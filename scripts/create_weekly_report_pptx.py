from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path('/Users/muse/Developer/muse-crm')
OUT = ROOT / '週報_MUSE_BBCRM_2026-05-25_05-29.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W, H = prs.slide_width, prs.slide_height

COLORS = {
    'bg': RGBColor(18, 18, 24),
    'panel': RGBColor(31, 31, 42),
    'purple': RGBColor(124, 58, 237),
    'violet': RGBColor(99, 102, 241),
    'cyan': RGBColor(6, 182, 212),
    'green': RGBColor(16, 185, 129),
    'amber': RGBColor(245, 158, 11),
    'red': RGBColor(239, 68, 68),
    'text': RGBColor(248, 250, 252),
    'muted': RGBColor(203, 213, 225),
    'dark_text': RGBColor(15, 23, 42),
    'light_bg': RGBColor(248, 250, 252),
    'light_panel': RGBColor(255, 255, 255),
    'border': RGBColor(226, 232, 240),
}

TITLE_FONT = 'Aptos Display'
BODY_FONT = 'Aptos'


def add_bg(slide, color='bg'):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    s.fill.solid(); s.fill.fore_color.rgb = COLORS[color]
    s.line.fill.background()
    return s


def textbox(slide, x, y, w, h, text, size: float = 18, color='text', bold=False, align='left', font=BODY_FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = COLORS[color]
    return tb


def add_title(slide, title, subtitle=None, dark=True):
    textbox(slide, 0.65, 0.35, 8.8, 0.55, title, 28, 'text' if dark else 'dark_text', True, font=TITLE_FONT)
    if subtitle:
        textbox(slide, 0.68, 0.93, 11.8, 0.35, subtitle, 11.5, 'muted' if dark else 'dark_text')


def card(slide, x, y, w, h, fill='panel', line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = COLORS[fill]
    if line:
        s.line.color.rgb = COLORS[line]; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def bullet_list(slide, x, y, w, h, items, size: float = 15, color='dark_text', bullet_color='purple'):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0; p.space_after = Pt(8)
        r = p.add_run(); r.text = '●  '; r.font.color.rgb = COLORS[bullet_color]; r.font.size = Pt(size)
        r2 = p.add_run(); r2.text = item; r2.font.name = BODY_FONT; r2.font.size = Pt(size); r2.font.color.rgb = COLORS[color]
    return tb


def add_image_fit(slide, path, x, y, w, h, border=True):
    path = str(ROOT / path)
    from PIL import Image
    im = Image.open(path)
    iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        width = w; height = w / img_ratio
        xx = x; yy = y + (h - height) / 2
    else:
        height = h; width = h * img_ratio
        xx = x + (w - width) / 2; yy = y
    pic = slide.shapes.add_picture(path, Inches(xx), Inches(yy), Inches(width), Inches(height))
    if border:
        pic.line.color.rgb = COLORS['border']; pic.line.width = Pt(1)
    return pic


def section_label(slide, x, y, text, fill='purple'):
    s = card(slide, x, y, 1.4, 0.32, fill=fill, radius=True)
    textbox(slide, x + 0.05, y + 0.055, 1.3, 0.22, text, 9.5, 'text', True, align='center')
    return s

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
card(slide, 0, 0, 13.333, 7.5, fill='bg', radius=False)
textbox(slide, 0.8, 0.78, 2.0, 0.35, 'MUSE BBCRM', 15, 'cyan', True)
textbox(slide, 0.8, 1.55, 9.8, 1.0, '本週修改與優化週報', 44, 'text', True, font=TITLE_FONT)
textbox(slide, 0.82, 2.55, 8.0, 0.5, '期間：2026/05/25 – 2026/05/29｜對象：主管與非技術團隊', 17, 'muted')
# hero cards
for i, (num, label, color) in enumerate([('51', '個主要檔案調整', 'purple'), ('3,110+', '行新增內容', 'cyan'), ('5/28', '完成多輪驗證', 'green'), ('PPTX', '含實際畫面截圖', 'amber')]):
    x = 0.85 + i * 3.05
    card(slide, x, 4.55, 2.65, 1.35, fill='panel', line=None)
    textbox(slide, x+0.2, 4.78, 2.1, 0.42, num, 24, color, True)
    textbox(slide, x+0.2, 5.25, 2.2, 0.38, label, 13, 'muted')
textbox(slide, 0.86, 6.65, 11.4, 0.32, '資料來源：本週專案工作區改動、QA 報告與 2026/05/28 實測截圖。', 10.5, 'muted')

# Slide 2 summary
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, 'light_bg'); add_title(slide, '一頁看懂本週成果', '重點放在使用者感受到的改善，不列複雜技術細節。', dark=False)
items = [
    ('對話工作台更完整', '客服可在同一畫面看對話、客戶狀態、AI 建議與快速回覆，操作更集中。', 'purple'),
    ('回覆效率提升', '預存語錄可搜尋、插入並帶入圖片，減少重複打字與回覆落差。', 'green'),
    ('主管視角更清楚', '儀表板補強角色檢查與資料範圍，降低一般人看到主管統計的風險。', 'cyan'),
    ('即時同步修補', '訊息送出、結案、待辦新增後，其他畫面可更快更新，減少手動刷新。', 'amber'),
]
for i,(h,b,c) in enumerate(items):
    x = 0.75 + (i%2)*6.05; y = 1.65 + (i//2)*2.1
    card(slide, x, y, 5.65, 1.55, fill='light_panel', line='border')
    card(slide, x+0.25, y+0.28, 0.45, 0.45, fill=c)
    textbox(slide, x+0.85, y+0.22, 4.35, 0.32, h, 18, 'dark_text', True)
    textbox(slide, x+0.85, y+0.67, 4.4, 0.55, b, 13, 'dark_text')
section_label(slide, 0.75, 6.15, '整體判斷', 'purple')
textbox(slide, 2.25, 6.12, 10.0, 0.5, '功能面已有明顯進展；但權限、AI 分析落地與庫存連線仍需收尾後再視為可正式上線。', 15, 'dark_text', True)

# Slide 3 screenshot inbox
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); add_title(slide, '實際畫面 1｜對話工作台整合', '畫面證據：客服能同時看到對話內容、客戶身份、AI 摘要與建議回覆。')
add_image_fit(slide, 'qa-output-17/ui-cockpit/inbox-cockpit.png', 0.55, 1.35, 8.65, 5.35)
card(slide, 9.45, 1.35, 3.25, 5.35, fill='panel')
textbox(slide, 9.75, 1.7, 2.65, 0.35, '本週改善重點', 18, 'text', True)
bullet_list(slide, 9.75, 2.25, 2.55, 3.9, [
    '把客戶狀態、意圖與摘要集中在右側，不必切換多個頁面。',
    '回覆草稿與建議動作更靠近輸入區，降低客服漏看機率。',
    '畫面已提供深色與淺色檢視，便於後續整理成正式操作介面。'
], size=13, color='muted', bullet_color='cyan')

# Slide 4 quick replies
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, 'light_bg'); add_title(slide, '實際畫面 2｜快速回覆與圖片輔助', '畫面證據：預存語錄彈窗可搜尋、分類，並能帶入圖片與標準回覆。', dark=False)
add_image_fit(slide, 'qa-output-17/tunnel-deploy/quick-replies-dropdown-open.png', 0.55, 1.35, 8.85, 5.35)
card(slide, 9.65, 1.35, 3.0, 5.35, fill='light_panel', line='border')
textbox(slide, 9.95, 1.72, 2.4, 0.35, '對業務的價值', 18, 'dark_text', True)
bullet_list(slide, 9.95, 2.25, 2.35, 3.8, [
    '常用話術集中管理，降低新人訓練成本。',
    '詢價、尺寸、材質說明可更快回覆。',
    '後續可擴充成「一鍵文字 + 圖片」完整話術包。'
], size=13, color='dark_text', bullet_color='purple')

# Slide 5 dashboard / permission
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); add_title(slide, '實際畫面 3｜主管儀表板與權限修補', '畫面證據：主管頁面已被列入角色檢查與 QA 範圍。')
add_image_fit(slide, 'qa-output-17/rbac-ui/admin-dashboard.png', 0.55, 1.35, 8.6, 5.35)
card(slide, 9.4, 1.35, 3.25, 5.35, fill='panel')
textbox(slide, 9.72, 1.72, 2.55, 0.35, '為什麼重要', 18, 'text', True)
bullet_list(slide, 9.72, 2.25, 2.55, 3.85, [
    '主管統計不應被一般帳號看到。',
    '本週已針對儀表板與側欄做修補。',
    '仍需再跑完整權限回歸，確認每種角色看到的內容正確。'
], size=13, color='muted', bullet_color='amber')

# Slide 6 AI analysis
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, 'light_bg'); add_title(slide, '實際畫面 4｜AI 客戶情報與語氣分析', '畫面證據：AI 區塊可顯示摘要、身份、急迫程度與回覆草稿；但仍需讓結果完全來自真實分析。', dark=False)
add_image_fit(slide, 'qa-output-17/screenshots-ai-tone/08-inbox-after-deep-analysis.png', 0.55, 1.35, 8.65, 5.35)
card(slide, 9.45, 1.35, 3.25, 5.35, fill='light_panel', line='border')
textbox(slide, 9.75, 1.72, 2.65, 0.35, '本週結論', 18, 'dark_text', True)
bullet_list(slide, 9.75, 2.25, 2.55, 3.85, [
    '介面入口與畫面呈現已具雛形。',
    '後台測試與前端建置通過。',
    '正式上線前，必須確認分析結果真的成功保存，不能只顯示暫時文字。'
], size=13, color='dark_text', bullet_color='red')

# Slide 7 quality gates
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); add_title(slide, '本週驗證結果', '以主管可讀方式整理：哪些已通過、哪些仍要補。')
rows = [
    ('已通過', '前端檢查與正式打包', '畫面可正常編譯，基本品質門檻通過。', 'green'),
    ('已通過', '後端重點測試', '快速回覆、互動訊息、即時更新等重點測試多數通過。', 'green'),
    ('需補強', '權限邊界', '曾發現一般帳號可看主管資料；已修補，需再跑完整驗證。', 'amber'),
    ('需補強', 'AI 分析保存', '操作入口可見，但分析結果未完全確認能寫入並回到畫面。', 'amber'),
    ('阻擋項', '庫存連線', '庫存頁曾遇到外部服務無法連線，正式上線前需處理。', 'red'),
]
for i,(status,title,desc,c) in enumerate(rows):
    y = 1.45 + i*0.92
    card(slide, 0.85, y, 11.65, 0.68, fill='panel')
    card(slide, 1.05, y+0.17, 1.15, 0.32, fill=c)
    textbox(slide, 1.12, y+0.205, 1.0, 0.2, status, 9, 'text', True, align='center')
    textbox(slide, 2.45, y+0.17, 2.3, 0.27, title, 15, 'text', True)
    textbox(slide, 5.0, y+0.17, 6.85, 0.3, desc, 13.2, 'muted')

# Slide 8 change areas
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, 'light_bg'); add_title(slide, '修改與優化分類', '本週不是單一功能，而是把對話、權限、AI、測試一起往產品化推進。', dark=False)
areas = [
    ('客服對話', ['送出訊息後即時更新', '支援回覆指定訊息', '可顯示機器人互動按鈕'], 'purple'),
    ('快速回覆', ['語錄搜尋與分類', '圖片附件基礎能力', '管理頁規劃擴充'], 'green'),
    ('主管與權限', ['主管統計加上角色檢查', '側欄依角色過濾', '降低資料誤看的風險'], 'cyan'),
    ('AI 客戶情報', ['接入新的分析來源選項', '語氣／急迫程度檢查', '補上成本與失敗風險提醒'], 'amber'),
    ('穩定性', ['即時通知修補', '測試報告與截圖留證', '部署前風險清單更清楚'], 'red'),
]
for i,(title,bullets,c) in enumerate(areas):
    x = 0.75 + (i%3)*4.1; y = 1.55 + (i//3)*2.35
    card(slide, x, y, 3.65, 1.8, fill='light_panel', line='border')
    card(slide, x+0.18, y+0.2, 0.35, 0.35, fill=c)
    textbox(slide, x+0.65, y+0.17, 2.6, 0.32, title, 16, 'dark_text', True)
    bullet_list(slide, x+0.28, y+0.65, 3.05, 0.9, bullets, size=10.8, color='dark_text', bullet_color=c)

# Slide 9 next steps
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); add_title(slide, '下週建議安排', '先收斂風險，再安排正式展示或部署。')
steps = [
    ('1', '權限完整回歸', '用管理員、主管、一般人三種身份重新驗證頁面與資料。'),
    ('2', 'AI 分析閉環', '確認點擊分析後，結果能保存、能回到畫面、失敗時有提示。'),
    ('3', '庫存服務連線', '確認庫存來源可用；若不可用，頁面要有友善提示。'),
    ('4', '快速回覆管理補齊', '補上圖片、類型、關鍵字管理，讓營運可自行維護話術。'),
    ('5', '整理上線包', '移除測試暫存檔與敏感本機檔案，再進入部署流程。'),
]
for i,(n,t,d) in enumerate(steps):
    x = 0.85 + (i%2)*6.0; y = 1.45 + (i//2)*1.5
    card(slide, x, y, 5.4, 1.0, fill='panel')
    card(slide, x+0.22, y+0.24, 0.52, 0.52, fill='purple')
    textbox(slide, x+0.38, y+0.33, 0.2, 0.2, n, 14, 'text', True, align='center')
    textbox(slide, x+0.95, y+0.18, 4.0, 0.28, t, 15.5, 'text', True)
    textbox(slide, x+0.95, y+0.53, 4.0, 0.25, d, 12.5, 'muted')

# Slide 10 appendix evidence
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, 'light_bg'); add_title(slide, '截圖與報告來源', '方便主管追溯本週證據。', dark=False)
refs = [
    'qa-output-17/ui-cockpit/inbox-cockpit.png',
    'qa-output-17/tunnel-deploy/quick-replies-dropdown-open.png',
    'qa-output-17/rbac-ui/admin-dashboard.png',
    'qa-output-17/screenshots-ai-tone/08-inbox-after-deep-analysis.png',
    'qa-output-17/QA_REPORT_17.md',
    'qa-output-17/RBAC_QA_REPORT_17.md',
    'qa-output-17/REALTIME_UPDATE_FIX_REPORT.md',
    'qa-output-17/AI_TONE_ANALYSIS_QA_REPORT.md',
]
for i, r in enumerate(refs):
    y = 1.45 + i*0.48
    textbox(slide, 0.9, y, 11.5, 0.25, f'• {r}', 13, 'dark_text')
textbox(slide, 0.9, 6.35, 11.4, 0.35, '備註：本週 Git 已提交紀錄為空；本週報依照目前工作區改動、5/28 QA 留證與實際畫面截圖統整。', 12, 'dark_text', True)

prs.save(str(OUT))
print(OUT)
