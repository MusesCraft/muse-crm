from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = Path('/Users/muse/Developer/muse-crm')
PPTX = ROOT / '週報_MUSE_BBCRM_2026-05-25_05-29.pptx'
BACKUP = ROOT / '週報_MUSE_BBCRM_2026-05-25_05-29.before_all_features.pptx'

COLORS = {
    'bg': RGBColor(18, 18, 24),
    'panel': RGBColor(31, 31, 42),
    'purple': RGBColor(124, 58, 237),
    'cyan': RGBColor(6, 182, 212),
    'green': RGBColor(16, 185, 129),
    'amber': RGBColor(245, 158, 11),
    'red': RGBColor(239, 68, 68),
    'text': RGBColor(248, 250, 252),
    'muted': RGBColor(203, 213, 225),
    'dark_text': RGBColor(15, 23, 42),
    'light_bg': RGBColor(248, 250, 252),
    'light_panel': RGBColor(255, 255, 255),
    'border': RGBColor(226, 232, 240),
}
TITLE_FONT = 'Aptos Display'
BODY_FONT = 'Aptos'

prs = Presentation(str(PPTX))
W, H = prs.slide_width, prs.slide_height

if not BACKUP.exists():
    BACKUP.write_bytes(PPTX.read_bytes())


def bg(slide, dark=True):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    s.fill.solid(); s.fill.fore_color.rgb = COLORS['bg' if dark else 'light_bg']
    s.line.fill.background()
    return s


def textbox(slide, x, y, w, h, text, size: float = 18, color='text', bold=False, align='left', font=BODY_FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = COLORS[color]
    return tb


def card(slide, x, y, w, h, fill='panel', line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = COLORS[fill]
    if line:
        s.line.color.rgb = COLORS[line]; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def add_title(slide, title, subtitle='', dark=True):
    textbox(slide, 0.55, 0.34, 10.6, 0.48, title, 25, 'text' if dark else 'dark_text', True, font=TITLE_FONT)
    if subtitle:
        textbox(slide, 0.58, 0.88, 11.7, 0.34, subtitle, 11.2, 'muted' if dark else 'dark_text')


def add_image_fit(slide, rel_path, x, y, w, h, border=True):
    path = ROOT / rel_path
    im = Image.open(path)
    iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        width = w; height = w / img_ratio
        xx = x; yy = y + (h - height) / 2
    else:
        height = h; width = h * img_ratio
        xx = x + (w - width) / 2; yy = y
    pic = slide.shapes.add_picture(str(path), Inches(xx), Inches(yy), Inches(width), Inches(height))
    if border:
        pic.line.color.rgb = COLORS['border']; pic.line.width = Pt(1)
    return pic

features = [
    ('登入頁', '帳號登入入口，含顯示密碼與忘記密碼提示。', 'qa-output-17/login-console-production.png', '使用者進入系統的第一步。'),
    ('主管儀表板', '集中查看客戶、對話、處理速度與待辦狀態。', 'qa-output-17/rbac-ui/admin-dashboard.png', '主管用來掌握整體營運狀況。'),
    ('對話收件匣', '客服主要工作台：對話列表、訊息內容、客戶情報同頁呈現。', 'qa-output-17/rbac-ui/admin-inbox.png', '日常客服回覆與追蹤的核心。'),
    ('對話工作台強化版', '右側整合客戶摘要、意圖與 AI 建議。', 'qa-output-17/ui-cockpit/inbox-cockpit.png', '減少客服切頁查資料。'),
    ('發送訊息', '直接在對話中回覆客戶並留下紀錄。', 'qa-output-17/screenshots/02-send.png', '確認客服回覆流程可操作。'),
    ('內部備註', '客服可在對話中留下只給內部看的備註。', 'qa-output-17/screenshots/03-internal-note.png', '避免內部提醒誤發給客戶。'),
    ('預存語錄', '快速插入常用回覆，並支援搜尋與分類。', 'qa-output-17/tunnel-deploy/quick-replies-dropdown-open.png', '提升回覆速度與一致性。'),
    ('求援／升級', '客服遇到需要主管協助的情況可求援。', 'qa-output-17/screenshots/05-escalation.png', '讓高風險或高價值客戶更快被主管看見。'),
    ('標記已解決', '完成服務後可將對話標記為已解決。', 'qa-output-17/screenshots/06-resolve.png', '利於後續追蹤處理進度。'),
    ('篩選與搜尋', '依狀態、渠道、關鍵字找到指定對話。', 'qa-output-17/screenshots/07-filters.png', '降低大量訊息中的查找成本。'),
    ('客戶管理', '查看客戶列表、來源與標籤。', 'qa-output-17/rbac-ui/admin-contacts.png', '把對話轉成可管理的客戶資料。'),
    ('客戶詳情', '查看單一客戶的基本資料、歷史與關聯資訊。', 'qa-output-17/screenshots/contact-王設計師.png', '支援更精準的跟進。'),
    ('報價管理', '查看與管理報價相關資訊。', 'qa-output-17/rbac-ui/admin-quotes.png', '銜接客服對話與商務報價。'),
    ('待辦事項', '集中管理需要後續處理的任務。', 'qa-output-17/rbac-ui/admin-actions.png', '確保追蹤事項不漏接。'),
    ('庫存管理', '查看商品與庫存相關資料。', 'qa-output-17/rbac-ui/admin-inventory.png', '讓客服與報價時能參考庫存狀態。'),
    ('知識庫', '集中維護產品、材質、話術等知識內容。', 'qa-output-17/rbac-ui/admin-knowledge-base.png', '讓 AI 與客服能有一致資訊來源。'),
    ('系統設定', '管理使用者、權限與系統相關設定。', 'qa-output-17/rbac-ui/admin-settings.png', '管理端維護入口。'),
    ('快速回覆管理', '設定頁可管理快捷回覆內容。', 'qa-output-17/tunnel-deploy/quick-replies-ui-smoke.png', '營運可維護常用話術。'),
    ('AI 客戶情報', '顯示客戶摘要、急迫程度與建議回覆。', 'qa-output-17/screenshots-ai-tone/08-inbox-after-deep-analysis.png', '協助客服判斷下一步。'),
    ('機器人互動訊息', '對話中可顯示 bot 按鈕與互動內容。', 'qa-output-17/bot-buttons-conversation-detail.png', '保留 Telegram/機器人互動脈絡。'),
    ('即時更新驗證', '不同畫面可接收新訊息更新，不必一直手動刷新。', 'qa-output-17/tunnel-deploy/realtime-api-to-inbox-smoke-final.png', '降低漏看新訊息的風險。'),
]

# Section divider
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide, dark=True)
textbox(slide, 0.85, 0.85, 2.0, 0.35, '功能截圖附錄', 15, 'cyan', True)
textbox(slide, 0.85, 1.55, 10.3, 0.85, '全功能畫面截圖', 40, 'text', True, font=TITLE_FONT)
textbox(slide, 0.88, 2.5, 10.8, 0.45, '以下將 BBCRM 主要功能逐一放入簡報，方便主管快速瀏覽目前產品覆蓋範圍。', 17, 'muted')
for i, (name, _, _, _) in enumerate(features):
    x = 0.9 + (i % 3) * 4.05
    y = 3.25 + (i // 3) * 0.48
    textbox(slide, x, y, 3.55, 0.25, f'{i+1:02d}. {name}', 11.2, 'muted')

# Feature slides
for idx, (name, desc, rel, value) in enumerate(features, start=1):
    dark = idx % 2 == 1
    slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide, dark=dark)
    add_title(slide, f'功能截圖 {idx:02d}｜{name}', desc, dark=dark)
    # image large on left, explanation card on right
    add_image_fit(slide, rel, 0.55, 1.35, 8.75, 5.35)
    fill = 'panel' if dark else 'light_panel'
    line = None if dark else 'border'
    card(slide, 9.55, 1.35, 3.0, 5.35, fill=fill, line=line)
    textbox(slide, 9.85, 1.72, 2.35, 0.32, '畫面說明', 17, 'text' if dark else 'dark_text', True)
    textbox(slide, 9.85, 2.22, 2.35, 0.9, desc, 13.2, 'muted' if dark else 'dark_text')
    textbox(slide, 9.85, 3.45, 2.35, 0.32, '對營運的價值', 17, 'text' if dark else 'dark_text', True)
    textbox(slide, 9.85, 3.95, 2.35, 0.95, value, 13.2, 'muted' if dark else 'dark_text')
    textbox(slide, 9.85, 5.95, 2.35, 0.28, f'截圖來源：{rel}', 7.8, 'muted' if dark else 'dark_text')

prs.save(str(PPTX))
print(f'updated={PPTX}')
print(f'backup={BACKUP}')
print(f'added_feature_slides={len(features)+1}')
print(f'total_slides={len(prs.slides)}')
